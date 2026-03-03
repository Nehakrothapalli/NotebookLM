import os, pathlib
import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from src.storage.paths import nb_root, ensure_tree
from src.storage.chroma_store import get_collection
from src.utils.text import safe_name

EMBED_MODEL = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")


def _file_path_from_gradio_obj(file_obj):
    if isinstance(file_obj, str):
        return file_obj
    path = getattr(file_obj, "name", None)
    if isinstance(path, str):
        return path
    return None

def simple_chunk(text: str, max_chars=2200, overlap=250):
    text = "\n".join([ln.strip() for ln in (text or "").splitlines() if ln.strip()]).strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]
    out, start = [], 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        out.append(text[start:end])
        if end == len(text): break
        start = max(0, end - overlap)
    return out

def extract_pdf(path: str):
    reader = PdfReader(path)
    items = []
    for i, page in enumerate(reader.pages):
        txt = (page.extract_text() or "").strip()
        if txt:
            items.append({"text": txt, "page": i+1})
    return items

def extract_pptx(path: str):
    prs = Presentation(path)
    items = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        txt = "\n".join(t.strip() for t in texts if t.strip()).strip()
        if txt:
            items.append({"text": txt, "slide": i+1})
    return items

def extract_txt(path: str):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        txt = f.read().strip()
    return [{"text": txt, "page": None}] if txt else []

def extract_url(url: str):
    r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
    r.raise_for_status()
    soup = BeautifulSoup(r.text, "html.parser")
    for tag in soup(["script","style","noscript"]):
        tag.decompose()
    text = soup.get_text("\n")
    text = "\n".join([ln.strip() for ln in text.splitlines() if ln.strip()])
    return [{"text": text[:200000], "page": None}]

def upsert_extracted(username: str, notebook_id: str, source_title: str, source_id: str, extracted_items: list[dict]) -> int:
    col = get_collection(username, notebook_id)
    ids, docs, metas = [], [], []
    chunk_counter = 0
    for item_idx, item in enumerate(extracted_items):
        for j, ch in enumerate(simple_chunk(item["text"])):
            ids.append(f"{source_id}::item{item_idx}::chunk{j}::{chunk_counter}")
            docs.append(ch)
            meta = {
                "source_title": source_title,
                "source_id": source_id,
                "page": item.get("page"),
                "slide": item.get("slide"),
            }
            meta = {k: v for k, v in meta.items() if v is not None}
            metas.append(meta)
            chunk_counter += 1
    if not docs:
        return 0
    embs = EMBED_MODEL.encode(docs, normalize_embeddings=True).tolist()
    col.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=embs)
    return len(docs)

def ingest_files(username: str, notebook_id: str, files) -> int:
    ensure_tree(username, notebook_id)
    raw_dir = os.path.join(nb_root(username, notebook_id), "files_raw")
    ex_dir = os.path.join(nb_root(username, notebook_id), "files_extracted")
    added = 0

    for f in (files or []):
        fp = _file_path_from_gradio_obj(f)
        if not fp:
            continue

        if not os.path.exists(fp):
            continue

        dest = os.path.join(raw_dir, os.path.basename(fp))
        pathlib.Path(dest).write_bytes(pathlib.Path(fp).read_bytes())

        ext = os.path.splitext(dest)[1].lower()
        if ext == ".pdf":
            extracted = extract_pdf(dest)
        elif ext == ".pptx":
            extracted = extract_pptx(dest)
        elif ext in [".txt", ".md"]:
            extracted = extract_txt(dest)
        else:
            continue

        ex_path = os.path.join(ex_dir, os.path.basename(dest) + ".txt")
        with open(ex_path, "w", encoding="utf-8") as ftxt:
            for item in extracted:
                loc = ""
                if item.get("page"):
                    loc = f"page={item.get('page')}"
                elif item.get("slide"):
                    loc = f"slide={item.get('slide')}"
                ftxt.write(f"\n--- {loc} ---\n{item['text']}\n")

        added += upsert_extracted(
            username,
            notebook_id,
            os.path.basename(dest),
            f"file:{os.path.basename(dest)}",
            extracted,
        )

    return added

def ingest_url(username: str, notebook_id: str, url: str) -> int:
    ensure_tree(username, notebook_id)
    extracted = extract_url(url)
    ex_dir = os.path.join(nb_root(username, notebook_id), "files_extracted")
    fname = safe_name(url.replace("https://","").replace("http://","").replace("/","_")) + ".txt"
    with open(os.path.join(ex_dir, fname), "w", encoding="utf-8") as f:
        f.write(extracted[0]["text"])
    return upsert_extracted(username, notebook_id, url, f"url:{url}", extracted)